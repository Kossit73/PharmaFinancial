import hashlib
import json
import os
from datetime import datetime, timezone
from typing import Any, Dict, List, Optional

import faiss
import pandas as pd
import pypdf
from docx import Document as DocxDocument
from fastapi import FastAPI, File, Form, UploadFile
from fastapi.responses import JSONResponse
from pydantic import BaseModel
from pptx import Presentation
from sentence_transformers import CrossEncoder, SentenceTransformer

EMBED_MODEL_NAME = os.getenv("EMBED_MODEL", "sentence-transformers/all-MiniLM-L6-v2")
RERANK_MODEL_NAME = os.getenv("RERANK_MODEL", "cross-encoder/ms-marco-MiniLM-L-6-v2")

CHUNK_TOKENS = 500
CHUNK_OVERLAP = 100
TOP_K = 12
RERANK_K = 5

DATA_DIR = os.getenv("DATA_DIR", "./projects")


class LLMClient:
    def __init__(self) -> None:
        self.provider = os.getenv("LLM_PROVIDER", "openai")

    def complete(self, prompt: str, system: Optional[str] = None) -> str:
        return "[LLM OUTPUT PLACEHOLDER]\n" + prompt[:500]


def sha256_file(path: str) -> str:
    digest = hashlib.sha256()
    with open(path, "rb") as handle:
        for chunk in iter(lambda: handle.read(8192), b""):
            digest.update(chunk)
    return digest.hexdigest()


def ensure_dirs(project_id: str) -> str:
    base = os.path.join(DATA_DIR, project_id)
    os.makedirs(os.path.join(base, "uploads"), exist_ok=True)
    os.makedirs(os.path.join(base, "parsed"), exist_ok=True)
    os.makedirs(os.path.join(base, "index"), exist_ok=True)
    os.makedirs(os.path.join(base, "financial"), exist_ok=True)
    return base


async def stream_save(upload: UploadFile, dest_path: str) -> None:
    with open(dest_path, "wb") as out:
        while True:
            chunk = await upload.read(1024 * 1024)
            if not chunk:
                break
            out.write(chunk)


def simple_tokenizer(text: str) -> List[str]:
    return text.split()


def chunk_text(text: str, chunk_tokens: int = CHUNK_TOKENS, overlap: int = CHUNK_OVERLAP) -> List[str]:
    tokens = simple_tokenizer(text)
    chunks: List[str] = []
    start = 0
    while start < len(tokens):
        end = min(len(tokens), start + chunk_tokens)
        chunk = " ".join(tokens[start:end])
        chunks.append(chunk)
        if end == len(tokens):
            break
        start = max(0, end - overlap)
    return chunks


def parse_pdf(path: str) -> List[Dict[str, Any]]:
    reader = pypdf.PdfReader(path)
    items: List[Dict[str, Any]] = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        items.append({"page_or_sheet": i, "text": text})
    return items


def parse_docx(path: str) -> List[Dict[str, Any]]:
    doc = DocxDocument(path)
    text = "\n".join(p.text for p in doc.paragraphs)
    return [{"page_or_sheet": 1, "text": text}]


def parse_pptx(path: str) -> List[Dict[str, Any]]:
    prs = Presentation(path)
    items: List[Dict[str, Any]] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        items.append({"page_or_sheet": i, "text": "\n".join(texts)})
    return items


def parse_txt(path: str) -> List[Dict[str, Any]]:
    with open(path, "r", encoding="utf-8", errors="ignore") as handle:
        return [{"page_or_sheet": 1, "text": handle.read()}]


EXCEL_KEYWORDS = ["NPV", "IRR", "DSCR", "Payback", "Capex", "Opex", "Revenue", "Discount", "Tax"]


def parse_excel_metrics(path: str) -> Dict[str, Any]:
    wb_hash = sha256_file(path)
    xls = pd.ExcelFile(path)

    snapshot: Dict[str, Any] = {
        "as_of": datetime.now(timezone.utc).isoformat(),
        "workbook_path": path,
        "workbook_hash": wb_hash,
        "currency": None,
        "assumptions": {},
        "capex_total": None,
        "opex_annual": None,
        "revenue_annual": None,
        "npv": None,
        "irr": None,
        "payback_years": None,
        "dscr_min": None,
        "sensitivities": [],
        "scenarios": [],
    }

    for sheet in xls.sheet_names:
        try:
            df = xls.parse(sheet)
        except Exception:
            continue
        df_str = df.astype(str)

        if snapshot["currency"] is None:
            if df_str.apply(
                lambda col: col.str.contains("USD|ZAR|EUR|\\$|R|€", case=False, regex=True, na=False)
            ).any().any():
                snapshot["currency"] = "DETECTED"

        for kw in EXCEL_KEYWORDS:
            matches = df_str.apply(
                lambda col: col.str.contains(fr"\\b{kw}\\b", case=False, regex=True, na=False)
            ).any(axis=1)
            if matches.any():
                row_idx = matches.idxmax()
                row = df.iloc[row_idx]
                val = None
                for v in row:
                    try:
                        val = float(str(v).replace(",", ""))
                        break
                    except Exception:
                        continue
                kw_l = kw.lower()
                if kw_l == "npv":
                    snapshot["npv"] = val
                elif kw_l == "irr":
                    snapshot["irr"] = val
                elif kw_l == "dscr":
                    snapshot["dscr_min"] = val
                elif kw_l == "payback":
                    snapshot["payback_years"] = val
                elif kw_l == "capex":
                    snapshot["capex_total"] = val
                elif kw_l == "opex":
                    snapshot["opex_annual"] = val
                elif kw_l == "revenue":
                    snapshot["revenue_annual"] = val
                elif kw_l == "discount":
                    snapshot["assumptions"].setdefault("discount_rate", val)
                elif kw_l == "tax":
                    snapshot["assumptions"].setdefault("tax_rate", val)

    return snapshot


class VectorIndex:
    def __init__(self, project_dir: str) -> None:
        self.project_dir = project_dir
        self.index_dir = os.path.join(project_dir, "index")
        self.index_path = os.path.join(self.index_dir, "faiss.index")
        self.meta_path = os.path.join(self.index_dir, "meta.jsonl")
        self.model = SentenceTransformer(EMBED_MODEL_NAME)
        self.dim = self.model.get_sentence_embedding_dimension()
        self.index = None
        self._load()

    def _load(self) -> None:
        if os.path.exists(self.index_path):
            self.index = faiss.read_index(self.index_path)
        else:
            self.index = faiss.IndexFlatIP(self.dim)
        if not os.path.exists(self.meta_path):
            open(self.meta_path, "a").close()

    def _save(self) -> None:
        faiss.write_index(self.index, self.index_path)

    def add(self, texts: List[str], metas: List[Dict[str, Any]]) -> None:
        embeddings = self.model.encode(texts, normalize_embeddings=True)
        self.index.add(embeddings)
        with open(self.meta_path, "a", encoding="utf-8") as handle:
            for meta, text in zip(metas, texts):
                handle.write(json.dumps({"meta": meta, "text": text}) + "\n")
        self._save()

    def search(self, query: str, top_k: int = TOP_K) -> List[Dict[str, Any]]:
        q = self.model.encode([query], normalize_embeddings=True)
        sims, idxs = self.index.search(q, top_k)
        results: List[Dict[str, Any]] = []
        with open(self.meta_path, "r", encoding="utf-8") as handle:
            lines = handle.readlines()
        for i in idxs[0]:
            if i < 0 or i >= len(lines):
                continue
            results.append(json.loads(lines[i]))
        return results


class Reranker:
    def __init__(self, name: str = RERANK_MODEL_NAME) -> None:
        try:
            self.model = CrossEncoder(name)
        except Exception:
            self.model = None

    def rerank(self, query: str, passages: List[Dict[str, Any]], k: int = RERANK_K) -> List[Dict[str, Any]]:
        if not self.model or not passages:
            return passages[:k]
        pairs = [(query, p["text"]) for p in passages]
        scores = self.model.predict(pairs)
        ranked = sorted(zip(passages, scores), key=lambda x: x[1], reverse=True)
        return [p for p, _ in ranked[:k]]


class IngestResponse(BaseModel):
    project_id: str
    files: List[str]


class GenerateRequest(BaseModel):
    project_id: str
    section_outline: Optional[List[str]] = None
    query_hint: Optional[str] = None


SYSTEM_PROMPT = (
    "You are a financial analyst producing a feasibility study. Only use facts from provided CONTEXT "
    "and FINANCIAL_SNAPSHOT. Cite sources inline like [Source: filename p.12] or [Sheet: Assumptions!B7]. "
    "If a claim is unsupported, say so. Keep each section structured, concise, and decision-oriented."
)

DEFAULT_SECTIONS = [
    "Executive Summary",
    "Project Description & Scope",
    "Market & Demand Analysis",
    "Technical & Operations",
    "Legal, Permitting & Environmental",
    "Implementation Plan",
    "Financial Analysis",
    "Risk Assessment & ESG",
    "Conclusion & Recommendation",
    "Appendices",
]

SECTION_PROMPT = """
[GOAL]
Draft the {section} for the feasibility study.

[GUIDANCE]
- Use FINANCIAL_SNAPSHOT metrics explicitly where relevant.
- Use CONTEXT passages with inline citations [Source: <file> p.<n>] or [Sheet: <name>!<cell>].
- State uncertainties and missing data.
- Avoid boilerplate; keep it specific to the project.

[FINANCIAL_SNAPSHOT]
{fin}

[CONTEXT]
{ctx}

[OUTPUT]
- 3–7 well-structured paragraphs
- Subheadings
- Bullet lists for key metrics & risks
- Citations inline
"""

app = FastAPI(title="RAG Feasibility Study Generator")
llm = LLMClient()


@app.post("/ingest", response_model=IngestResponse)
async def ingest(project_id: str = Form(...), files: List[UploadFile] = File(...)) -> IngestResponse:
    base = ensure_dirs(project_id)
    saved: List[str] = []
    for file in files:
        dest = os.path.join(base, "uploads", file.filename)
        await stream_save(file, dest)
        saved.append(dest)

        parsed_items: List[Dict[str, Any]] = []
        ext = os.path.splitext(file.filename)[1].lower()
        if ext == ".pdf":
            parsed_items = parse_pdf(dest)
        elif ext == ".docx":
            parsed_items = parse_docx(dest)
        elif ext == ".pptx":
            parsed_items = parse_pptx(dest)
        elif ext in (".txt", ".md"):
            parsed_items = parse_txt(dest)
        elif ext in (".xlsx", ".xls", ".csv"):
            try:
                if ext == ".csv":
                    df = pd.read_csv(dest)
                    text = df.to_csv(index=False)
                    parsed_items = [{"page_or_sheet": 1, "text": text}]
                else:
                    xls = pd.ExcelFile(dest)
                    for sheet in xls.sheet_names:
                        df = xls.parse(sheet)
                        text = df.to_csv(index=False)
                        parsed_items.append({"page_or_sheet": sheet, "text": text})
            except Exception:
                parsed_items = []

        parsed_path = os.path.join(base, "parsed", f"{os.path.basename(dest)}.jsonl")
        with open(parsed_path, "w", encoding="utf-8") as out:
            for item in parsed_items:
                out.write(json.dumps(item) + "\n")

        vi = VectorIndex(base)
        texts: List[str] = []
        metas: List[Dict[str, Any]] = []
        for item in parsed_items:
            chunks = chunk_text(item.get("text", ""))
            for chunk in chunks:
                texts.append(chunk)
                metas.append(
                    {
                        "project_id": project_id,
                        "file_path": dest,
                        "file_type": ext[1:],
                        "page_or_sheet": item.get("page_or_sheet"),
                        "section": None,
                        "char_start": 0,
                        "char_end": len(chunk),
                        "hash": sha256_file(dest),
                    }
                )
        if texts:
            vi.add(texts, metas)

        if ext in (".xlsx", ".xls"):
            snap = parse_excel_metrics(dest)
            with open(os.path.join(base, "financial", "snapshot.json"), "w", encoding="utf-8") as sf:
                json.dump(snap, sf, indent=2)

    return IngestResponse(project_id=project_id, files=saved)


@app.post("/generate")
def generate(req: GenerateRequest) -> JSONResponse:
    base = ensure_dirs(req.project_id)
    vi = VectorIndex(base)
    reranker = Reranker()

    snap_path = os.path.join(base, "financial", "snapshot.json")
    financial: Dict[str, Any] = {}
    if os.path.exists(snap_path):
        with open(snap_path, "r", encoding="utf-8") as handle:
            financial = json.load(handle)

    hint = req.query_hint or "feasibility study for project"
    passages = vi.search(hint, top_k=TOP_K)
    passages = reranker.rerank(hint, passages, k=RERANK_K)

    def fmt_passage(passage: Dict[str, Any]) -> str:
        meta = passage.get("meta", {})
        fname = os.path.basename(meta.get("file_path", "source"))
        page = meta.get("page_or_sheet", "?")
        return f"[Source: {fname} p.{page}]\n" + passage.get("text", "")

    ctx_block = "\n\n".join(fmt_passage(p) for p in passages)

    sections = req.section_outline or DEFAULT_SECTIONS

    fin_lines: List[str] = []
    for key, value in financial.items():
        if isinstance(value, (str, int, float)):
            fin_lines.append(f"- {key}: {value}")
    if financial.get("assumptions"):
        fin_lines.append("- assumptions:")
        for ak, av in financial["assumptions"].items():
            fin_lines.append(f"  - {ak}: {av}")
    fin_block = "\n".join(fin_lines) or "- (no financial snapshot found)"

    outputs: List[Dict[str, str]] = []
    for sec in sections:
        prompt = SECTION_PROMPT.format(section=sec, fin=fin_block, ctx=ctx_block)
        txt = llm.complete(prompt, system=SYSTEM_PROMPT)
        outputs.append({"section": sec, "content": txt})

    report = {
        "project_id": req.project_id,
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "sections": outputs,
        "financial_snapshot": financial,
    }

    out_json = os.path.join(base, "report.json")
    out_md = os.path.join(base, "report.md")

    with open(out_json, "w", encoding="utf-8") as jf:
        json.dump(report, jf, indent=2)

    with open(out_md, "w", encoding="utf-8") as mf:
        mf.write("# Feasibility Study\n\n")
        mf.write(f"**Project:** {req.project_id}\n\n")
        mf.write(f"**Generated at:** {report['generated_at']}\n\n")
        for sec in outputs:
            mf.write(f"## {sec['section']}\n\n{sec['content']}\n\n")
        mf.write("---\n\n### Financial Snapshot\n\n")
        mf.write("```json\n" + json.dumps(financial, indent=2) + "\n```\n")

    return JSONResponse(
        {
            "project_id": req.project_id,
            "report_json": out_json,
            "report_md": out_md,
            "sections": [s["section"] for s in outputs],
        }
    )


if __name__ == "__main__":
    import uvicorn

    uvicorn.run(app, host="0.0.0.0", port=8000)
